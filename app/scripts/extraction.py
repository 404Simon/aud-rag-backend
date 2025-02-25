import re

import PyPDF2
from pptx import Presentation


def extract_text_from_pdf(pdf_path: str) -> list[tuple[int, str]]:
    with open(pdf_path, "rb") as file:
        reader = PyPDF2.PdfReader(file)
        pages: list[tuple[int, str]] = []
        for page_num, page in enumerate(reader.pages, 1):
            text = filter_page_content(page.extract_text())
            if text:
                pages.append((page_num, text))
    return pages


def extract_text_from_pptx(pptx_path: str) -> list[tuple[int, str]]:
    presentation = Presentation(pptx_path)
    pages: list[tuple[int, str]] = []

    for idx, slide in enumerate(presentation.slides, 1):
        slide_content = []

        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())

            # Extract text from tables
            if shape.has_table:
                table_text = []
                for row in shape.table.rows:
                    row_text = [
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    ]
                    if row_text:
                        table_text.append(" | ".join(row_text))
                if table_text:
                    slide_content.append("\nTable content:\n" + "\n".join(table_text))

        if slide_content:
            pages.append((idx, "\n".join(slide_content).strip()))

    return pages


def filter_page_content(text: str) -> str | None:
    # remove everything until a date like 02/2022 with regex
    text = re.sub(r".*?(\d{2}/\d{4})", "", text, flags=re.DOTALL).strip()

    # remove header adopted since 7th lecture
    text = re.sub(
        r"""Algorithmen und Datenstrukturen\s*
DHBW Stuttgart Campus Horb\s*\d+\s*-\s*\d+""",
        "",
        text,
        flags=re.DOTALL,
    ).strip()

    text = re.sub(r"^\d+\s*-\s*\d+\s*,?\s*", "", text).strip()

    text = text.replace("\x00", "")

    if text and "Übung" not in text:
        return text.strip()
    return None


def extract_single_page_from_pdf(pdf_path: str, page_num: int) -> str:
    with open(pdf_path, "rb") as file:
        reader = PyPDF2.PdfReader(file)
        page = reader.pages[page_num]
        text = page.extract_text()
        text = filter_page_content(text)
    return text
